
from io import BytesIO
import pandas as pd
import openai
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.util import Inches
import os
import streamlit as st
import pandas as pd
import openai
from dotenv import load_dotenv
from openai import OpenAI
load_dotenv()

st.title("Data Report Generator 📈")

# def check_password():
#     """Returns `True` if the user had a correct password."""

#     def login_form():
#         """Form with widgets to collect user information"""
#         with st.form("Credentials"):
#             st.text_input("Username", key="username")
#             st.text_input("Password", type="password", key="password")
#             st.form_submit_button("Log in", on_click=password_entered)

#     def password_entered():
#         """Checks whether a password entered by the user is correct."""
#         if st.session_state["username"] in st.secrets[
#             "passwords"
#         ] and hmac.compare_digest(
#             st.session_state["password"],
#             st.secrets.passwords[st.session_state["username"]],
#         ):
#             st.session_state["password_correct"] = True
#             del st.session_state["password"]  # Don't store the username or password.
#             del st.session_state["username"]
#         else:
#             st.session_state["password_correct"] = False

#     # Return True if the username + password is validated.
#     if st.session_state.get("password_correct", False):
#         return True

#     # Show inputs for username + password.
#     login_form()
#     if "password_correct" in st.session_state:
#         st.error("😕 User not known or password incorrect")
#     return False

# if not check_password():
#     st.stop()

# Importing the datasets
sales_Data_name = st.file_uploader("Upload your Data File in XLSX format", type=['xlsx'])
target_Data_name = st.file_uploader("Upload your Target file in XLSX format", type=['xlsx'])
Current_year_Selection = st.selectbox("Select Current Year",[2023, 2024])
Current_month_Selection = st.selectbox("Select Current Month",range(1,13))
Previous_year_Selection = Current_year_Selection-1
Previous_month_selection = Current_month_Selection
YTD_Month=10

if st.button("Generate Report"):
    with st.spinner("Loading .. 🔃"):
    #   #load the presentation template
        presentation = Presentation("Sales Deck - Output.pptx")

        if sales_Data_name is not None and target_Data_name is not None:
            try:
                sales_data = pd.read_excel(sales_Data_name)
                target_data = pd.read_excel(target_Data_name)
            except Exception as e:
                st.write("Error reading excel files - ", e)  
        st.title("Data Report ")
        # Display the first few rows of each dataset to understand their structure
     





        # Convert the 'Date' columns to datetime
        sales_data['Date'] = pd.to_datetime(sales_data['Date'])
        target_data['Date'] = pd.to_datetime(target_data['Date'])


        # Filter sales data
        sales= sales_data[(sales_data['Date'].dt.year == Current_year_Selection) & (sales_data['Date'].dt.month == Current_month_Selection)]


        previous_sales= sales_data[(sales_data['Date'].dt.year == Previous_year_Selection) & (sales_data['Date'].dt.month == Previous_month_selection)]


        adjusted_year = Current_year_Selection if Current_month_Selection >= 10 else Current_year_Selection - 1

        # Construct the mask for filtering the data
        mask = (
            ((sales_data['Date'].dt.year == adjusted_year) & (sales_data['Date'].dt.month >= 10)) |  # Includes from October of the starting YTD year
            ((sales_data['Date'].dt.year == Current_year_Selection) & (sales_data['Date'].dt.month <= Current_month_Selection))  # Includes up to the selected month of the current year
        )

        # Apply the mask to filter the DataFrame
        Current_year_to_date_sales = sales_data[mask]


        # Adjust the year based on whether the previous month selection is less than the starting YTD month (October)
        YTD_start_year = Previous_year_Selection if Previous_month_selection >= 10 else Previous_year_Selection - 1

        # Construct the mask for filtering the data
        mask = (
            ((sales_data['Date'].dt.year == YTD_start_year) & (sales_data['Date'].dt.month >= 10)) |  # Includes from October of the starting YTD year
            ((sales_data['Date'].dt.year == Previous_year_Selection) & (sales_data['Date'].dt.month <= Previous_month_selection))  # Includes up to the selected month of the current year
        )

        # Apply the mask to filter the DataFrame
        Previous_year_to_date_sales = sales_data[mask]


        print(Current_year_to_date_sales)


        print(Previous_year_to_date_sales)


        print(previous_sales)


        # Calculate total sales
        total_sales = round(sales['Sales_FC'].sum())


        total_previous_sales = round(previous_sales['Sales_FC'].sum())


        current_year_to_date_sales = round(Current_year_to_date_sales['Sales_FC'].sum())


        previous_year_to_date_sales = round(Previous_year_to_date_sales['Sales_FC'].sum())


        print(total_sales)


        print(total_previous_sales)


        print(current_year_to_date_sales)


        print(previous_year_to_date_sales)


        growth_yoy = round((total_sales - total_previous_sales)/total_previous_sales,1)
        print(growth_yoy)


        growth_year_to_date = round((current_year_to_date_sales-previous_year_to_date_sales)/previous_year_to_date_sales,1)
        print(growth_year_to_date)


        # Filter target data for November 2023 and get the Company Total
        target= target_data[(target_data['Date'].dt.year == Current_year_Selection) & (target_data['Date'].dt.month == Current_month_Selection)]
        company_total_target = round(target[target['Attributes'] == 'Company Total']['Value'].sum())


        print(company_total_target)


        difference_sales_target =round(total_sales-company_total_target)


        print(difference_sales_target )


        sales_to_target_ratio = round((total_sales / company_total_target)*100,1)



        print(sales_to_target_ratio)


        # Store the variables
        data = {
            'Total Sales': total_sales,
            'Company Total Target': company_total_target,
            'Difference (Sales - Target)': difference_sales_target ,
            'Sales/Target Ratio': sales_to_target_ratio,
            'Total Previous Sales': total_previous_sales,
            'Growth':growth_yoy,
            'Current Year to Date Sales':current_year_to_date_sales,
            'Previous Year to Date Sales':previous_year_to_date_sales,
            'Growth Year to Date':growth_year_to_date
        }


        print(data)


        prompt = f"""
        Here is the sales data for the Period:
        Total Sales: {data['Total Sales']}
        Company Total Target: {data['Company Total Target']}
        Difference (Sales - Target): {data['Difference (Sales - Target)']}
        Sales to target ratio :{data['Sales/Target Ratio']}


        Please provide a brief two line observation based on the summary of this data and include the difference of Sales - target into the observation
        """


        client= openai.OpenAI(api_key=os.getenv("OPEN_API_KEY"))


        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst with ability to generate quality observation from the given data"},
            {"role": "user", "content": prompt}
        ]
        )


        # Print the observation
        observation=completion.choices[0].message.content
        print(observation)


        prompt2 = f"""
        Here is the sales data for the current year and year on year:
        Current Month Sales: {total_sales}
        Year on Year Sales: {total_previous_sales}
        Growth/Decline from year on year sales to current month sales: {growth_yoy}


        Please provide a brief two line observation based on the summary of this data.
        """


        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst with ability to generate quality observation from the given data"},
            {"role": "user", "content": prompt2}
        ]
        )


        # Print the observation
        observation2=completion.choices[0].message.content
        print(observation2)


        prompt3 = f"""
        Here is the current year to date sales, previous year to date sales and growth/decline of previous to current year to date sales :
        Current year to date sales: {data['Current Year to Date Sales']}
        Preivous year to date sales: {data['Previous Year to Date Sales']}
        Growth/Decline of previous to current year to date sales: {data['Growth Year to Date']}


        Please provide a brief two line observation based on the summary of this data.
        """


        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst with ability to generate quality observation from the given data"},
            {"role": "user", "content": prompt3}
        ]
        )


        # Print the observation
        observation3=completion.choices[0].message.content
        print(observation3)


        prompt4 = f"""
        Here are the three observations :
        Target preformance: {observation}
        Year on Year sales performance: {observation2}
        Year to Date Performance: {observation3}


        Please provide a two line observation in one single para based on the three observations. Accuracy is important and read the observations correctly.
        """



        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst with ability to generate quality observation from the given data"},
            {"role": "user", "content": prompt4}
        ]
        )


        # Print the observation
        observation4=completion.choices[0].message.content
        print(observation4)


        df = pd.DataFrame(list(data.items()), columns=['Metric', 'Value'])


        print(df)


        chart_data = CategoryChartData()
        chart_data.categories = ['Category 1', 'Category 2']
        chart_data.add_series('Series 1', (sales_to_target_ratio, 0 if sales_to_target_ratio >= 1 else 1 - sales_to_target_ratio ))


        chart_data2 = CategoryChartData()
        chart_data2.categories = ['Category 1', 'Category 2']
        chart_data2.add_series('PY', [total_previous_sales])
        chart_data2.add_series('CY',[total_sales])


        chart_data3 = CategoryChartData()
        chart_data3.categories = ['Category 1', 'Category 2']
        chart_data3.add_series('PYTD', [previous_year_to_date_sales])
        chart_data3.add_series('CYTD',[current_year_to_date_sales])


        def update_chart_data(slide, chart_shape_index, chart_data):
            chart = slide.shapes[chart_shape_index].chart
            chart.replace_data(chart_data)



        def update_chart_data(slide, chart_shape_index, chart_data2):
            chart = slide.shapes[chart_shape_index].chart
            chart.replace_data(chart_data2)



        def update_chart_data(slide, chart_shape_index, chart_data3):
            chart = slide.shapes[chart_shape_index].chart
            chart.replace_data(chart_data3)



        # Update the chart data on the desired slide
        slide_index = 1  # Index of the slide containing the chart (0 for the first slide)
        chart_shape_index = 10  # Index of the chart shape on the slide (adjust based on your slide layout)


        #Update the chart data on the desired slide
        #Index of the slide containing the chart (0 for the first slide)
        chart_shape_index2 = 11  # Index of the chart shape on the slide (adjust based on your slide layout)


        # Update the chart data on the desired slide
        # Index of the slide containing the chart (0 for the first slide)
        chart_shape_index3 = 20  # Index of the chart shape on the slide (adjust based on your slide layout)


        slide = presentation.slides[slide_index]
        update_chart_data(slide, chart_shape_index, chart_data)


        slide = presentation.slides[slide_index]
        update_chart_data(slide, chart_shape_index2, chart_data2)


        slide = presentation.slides[slide_index]
        update_chart_data(slide, chart_shape_index3, chart_data3)


        # Select the second slide (index 1 because it's zero-indexed)
        slide = presentation.slides[1]

        # Dictionary of index: new_text
        text_updates = {
            1: f"The total sales units for the Current Period reached {total_sales}",
            2: observation,
            3: observation2,
            4: observation3,
            5: observation4
        }



        # Iterate through each update instruction
        for target_index, new_text in text_updates.items():
            textbox_counter = 0  # Reset counter for each update
            # Iterate through the shapes in the slide
            for shape in slide.shapes:
                # Check if the shape is a text box
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    # If the counter matches the target index, update the text
                    if textbox_counter == target_index:
                        if shape.has_text_frame:
                            # To retain formatting, replace the text of each run
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = new_text
                        break
                    textbox_counter += 1


        # Filter the dataset for sales in November 2022 and November 2023
        sales_data['Year'] = sales_data['Date'].dt.year
        sales_data['Month'] = sales_data['Date'].dt.month

        # Filter sales data for November of the previous year
        previous_periodsales= sales_data[(sales_data['Year'] == Previous_year_Selection) & 
                                        (sales_data['Month'] == Previous_month_selection)]

        # Filter sales data for November of the current year
        current_periodsales= sales_data[(sales_data['Year'] == Current_year_Selection) & 
                                        (sales_data['Month'] == Current_month_Selection)]




        # Calculate the total sales for November 2022 and November 2023 by Attributes
        PP_sales_by_attributes = previous_periodsales.groupby('Attributes')['Sales_FC'].sum()
        CP_sales_by_attributes = current_periodsales.groupby('Attributes')['Sales_FC'].sum()

        # Merge the sales data for both years
        sales_comparison = pd.DataFrame({
            'Previous Period': PP_sales_by_attributes,
            'Current Period': CP_sales_by_attributes
        }).fillna(0)

        # Calculate the change in sales
        sales_comparison['Change'] = sales_comparison['Current Period'] - sales_comparison['Previous Period']


        fsj_tank_total_change = sales_comparison.at['FSJ Tank Total', 'Change']
        fsj_vac_total_change = sales_comparison.at['FSJ Vac Total', 'Change']
        fort_nelson_total_change = sales_comparison.at['Fort Nelson Total', 'Change']


        data = {
            'Total Sales': total_sales,
            'FSJ Tank': fsj_tank_total_change,
            'FSJ Vac': fsj_vac_total_change,
            'FN': fort_nelson_total_change,
            'Total Previous Sales': total_previous_sales,
        }
        print (data)


        data= {'category': ['Previous Total','FSJ Tank', 'FSJ Vac', 'FN', 'Total Sales'],
                'value': [total_previous_sales, fsj_tank_total_change, fsj_vac_total_change, fort_nelson_total_change, total_sales]}


        print(data)


        # Data provided by the user
        data= {'category': ['Previous Total','FSJ Tank', 'FSJ Vac', 'FN', 'Total Sales'],
                'value': [total_previous_sales, fsj_tank_total_change, fsj_vac_total_change, fort_nelson_total_change, total_sales]}

        # Creating DataFrame
        df = pd.DataFrame(data)

        # Adjusted cumulative values for plotting
        cumulative_shifted = df['value'].shift(1).fillna(0).cumsum()

        # Plotting the waterfall chart
        fig, ax = plt.subplots()

        # Set the background color
        fig.patch.set_facecolor('#F6F6F7')
        ax.set_facecolor('#F6F6F7')

        # Plot bars with colors based on value changes
        for i in range(len(df)):
            if df['category'][i] == 'Total Sales':
                bottom = 0
                color = '#717EEE'
            elif df['category'][i] == 'Previous Total':
                bottom = 0
                color = '#FFAA60'
            else:
                bottom = cumulative_shifted[i]
                color = '#50D0BC' if df['value'][i] >= 0 else '#EA738D'
            ax.bar(df['category'][i], df['value'][i], bottom=bottom, color=color)
            
            # Adding data labels with currency formatting and rounding
            height = bottom + df['value'][i]
            label = f'${df["value"][i]:,.0f}'
            ax.text(i, height, label, ha='center', va='bottom' if df['value'][i] >= 0 else 'top')

        # Adding a horizontal line for the previous total
        ax.axhline(y=total_previous_sales, color='grey', linewidth=0.8, linestyle='--')

        # Remove grid, Y-axis, and Y-axis line
        ax.grid(False)
        ax.yaxis.set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['top'].set_visible(False)

        # Remove Y axis label, Y axis title, X axis title, and chart title
        ax.yaxis.set_label_text('')
        ax.set_ylabel('')
        ax.set_xlabel('')
        plt.title('')

        # Labels and rotation
        plt.xticks(rotation=0)

        # Display the plot
        plt.show()

        img_stream = BytesIO()
        fig.savefig(img_stream, format='png')
        img_stream.seek(0)

        # Close the plot
        plt.close(fig)

        slide = presentation.slides[2]

        # Define the position and size of the image
        left = Inches(5.5551181)
        top = Inches(1.740157)
        width = Inches(6.4015748)
        height = Inches(4.7992126)

        # Insert the image
        slide.shapes.add_picture(img_stream, left, top, width, height)
        img_stream.truncate(0)  # Truncate the buffer to zero bytes





        prompt5 = f"""
        Here is the data which has Previous total and Current Total Sales and breakdown of change between the two year by Location FSJ Tank, FSJ Vac and FN
            'Total Previous Sales': {[total_previous_sales]},
            'FSJ Tank': {[fsj_tank_total_change]},
            'FSJ Vac': {[fsj_vac_total_change]},
            'FN': {[fort_nelson_total_change]},
            'Total Sales': {[total_sales]}

        Please provide observation in one single para based on the data and given an insight on contribution to growth or decline. Write the observation to make to probe the data further.
        """


        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst with ability to generate quality observation from the given data"},
            {"role": "user", "content": prompt5}
        ]
        )


        # Print the observation
        observation5=completion.choices[0].message.content
        print(observation5)


        # Select the second slide (index 1 because it's zero-indexed)
        slide = presentation.slides[2]

        # Dictionary of index: new_text
        text_updates = {
            1: observation5,
            
        }


        # Iterate through each update instruction
        for target_index, new_text in text_updates.items():
            textbox_counter = 0  # Reset counter for each update
            # Iterate through the shapes in the slide
            for shape in slide.shapes:
                # Check if the shape is a text box
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    # If the counter matches the target index, update the text
                    if textbox_counter == target_index:
                        if shape.has_text_frame:
                            # To retain formatting, replace the text of each run
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = new_text
                        break
                    textbox_counter += 1




        import tempfile

        # Group by 'Attributes' and 'work_type' and sum the sales for each group
        PP_Sales = previous_periodsales.groupby(['Attributes', 'work_type'])['Sales_FC'].sum().reset_index()
        CP_Sales =current_periodsales.groupby(['Attributes', 'work_type'])['Sales_FC'].sum().reset_index()

        # Merge the two dataframes on 'Attributes' and 'work_type'
        merged_df = pd.merge(PP_Sales, CP_Sales, on=['Attributes', 'work_type'], suffixes=('_PP', '_CP'))

        # Calculate the change in sales
        merged_df['Change'] = merged_df['Sales_FC_CP'] - merged_df['Sales_FC_PP']

        # Sort the dataframe by 'Change' in descending order for plotting
        waterfall_data_sorted = merged_df.sort_values(by='Change', ascending=False)

        # Function to remove 'Total' from the attribute values and shorten the labels
        def remove_total_and_shorten(label):
            parts = label.split(' - ')
            attribute = parts[0].replace('Total', '').strip()
            attribute = parts[0].replace('Fort', '').strip()
            work_type = parts[1]
            shortened = attribute[:4] + ' - ' + work_type[:5]  # Taking first 5 characters of each part
            return shortened

        # Apply the label modification
        modified_labels = waterfall_data_sorted.apply(lambda row: remove_total_and_shorten(row['Attributes'] + ' - ' + row['work_type']), axis=1)

        # Plotting the waterfall chart with more spacing and modified labels
        def img_stream_waterfall(df, short_labels, title, presentation, slide_index):
            fig, ax = plt.subplots(figsize=(12.42126, 4.7125984))

            # Initialize values
            start = total_previous_sales
            end = total_sales
            changes = df['Change'].values
            colors = ['#50D0BC' if x >= 0 else '#EA738D' for x in changes]

            # Plot initial value
            ax.bar('Previous Period', start, color='#FFAA60')
            ax.text(0, start + (start * 0.04), f'{start:.2f}', ha='center', va='bottom')  # Positioning start value below the bar
            
            # Plot changes with more spacing
            running_total = start
            for i in range(len(changes)):
                ax.bar(i + 1, changes[i], bottom=running_total, color=colors[i], width=1)
                running_total += changes[i]

            # Add data labels above the bars
            running_total = start
            for i in range(len(changes)):
                height = changes[i]
                ax.text(i + 1, running_total + (height * 0.5), f'{height:.2f}', ha='center', va='top')  # Positioning labels above the bars
                running_total += height

            # Plot final value
            ax.bar(len(changes) + 1, end, color='#717EEE')
            ax.text(len(changes) + 1, end + (end * 0.04), f'{end:.2f}', ha='center', va='top')  # Positioning final value above the bar

            

            # Formatting
            ax.set_xticks(range(len(changes) + 2))
            ax.set_xticklabels(['Previous Period'] + short_labels.tolist() + ['Current Period'], rotation=0)
            plt.ylabel('Sales')
            plt.title(title)
            plt.grid(axis='y')
            ax.yaxis.set_visible(False)
            ax.grid(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)

            # Set the background color
            fig.patch.set_facecolor('#F6F6F7')
            ax.set_facecolor('#F6F6F7')

            plt.tight_layout()  # Ensure everything fits within the figure area
            plt.show()
            
            # Save the figure directly to a temporary file
            temp_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            fig.savefig(temp_file.name, format='png')

        # Get reference to the slide
            slide = presentation.slides[slide_index]

        # Define the position and size of the image
            left = Inches(0.354331)
            top = Inches(0.984252)
            width = Inches(12.42126)
            height = Inches(4.7125984)

        # Insert the image from the temporary file
            slide.shapes.add_picture(temp_file.name, left, top, width, height)

        # Close the figure and remove the temporary file
            plt.close(fig)
            temp_file.close()

            

        # Plotting the waterfall chart with more spacing and modified labels
        img_stream_waterfall(waterfall_data_sorted, modified_labels,'', presentation, 3)


        # Plotting the waterfall chart with more spacing and modified labels
        img_stream_waterfall = BytesIO()
        fig.savefig(img_stream_waterfall, format='png')
        img_stream_waterfall.seek(0)

        # Close the plot
        plt.close(fig)

        # slide = presentation.slides[3]

        # # Define the position and size of the image
        # left = Inches(5.5551181)
        # top = Inches(1.740157)
        # width = Inches(6.4015748)
        # height = Inches(4.7992126)

        # # Insert the image
        # slide.shapes.add_picture(img_stream_waterfall, left, top, width, height)

        # # Clear and seek the BytesIO object
        # img_stream_waterfall.truncate(0)
        # img_stream_waterfall.seek(0)





        prompt6 = f"""
        Here is the data which has Previous total and Current Total Sales and breakdown of change between the two year by Location and work type
            Change:{[waterfall_data_sorted]}
            'Total Sales': {[total_sales]}

        Please provide observation in one single para in 150 words based on the data and given an insight on contribution to growth or decline. Write the observation to make to probe the data further.
        """


        completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst with ability to generate quality observation from the given data"},
            {"role": "user", "content": prompt6}
        ]
        )


        # Print the observation
        observation6=completion.choices[0].message.content
        print(observation6)


        # Select the second slide (index 1 because it's zero-indexed)
        slide = presentation.slides[3]

        # Dictionary of index: new_text
        text_updates1 = {
            2: observation6,
        }



        # Iterate through each update instruction
        for target_index, new_text in text_updates1.items():
            textbox_counter = 0  # Reset counter for each update
            # Iterate through the shapes in the slide
            for shape in slide.shapes:
                # Check if the shape is a text box
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    # If the counter matches the target index, update the text
                    if textbox_counter == target_index:
                        if shape.has_text_frame:
                            # To retain formatting, replace the text of each run
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = new_text
                        break
                    textbox_counter += 1
                    


        # Your existing code here

        presentation.save("Sales Deck - Output.pptx")
        with open(r"Sales Deck - Output.pptx", "rb") as file:
            btn = st.download_button(
            label="Download",
            data=file,
            file_name="report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
          )
